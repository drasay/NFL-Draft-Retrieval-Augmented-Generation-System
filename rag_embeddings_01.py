import os
import fitz 
from pptx import Presentation
import chromadb
from sentence_transformers import SentenceTransformer

# Setup: Paths
DOC_DIR = "C:/python/defensive_line_prospects"
raw_doc_db_path = os.path.join(DOC_DIR, "chroma_raw_documents")
embedding_db_minilm = os.path.join(DOC_DIR, "chroma_minilm")
embedding_db_e5 = os.path.join(DOC_DIR, "chroma_e5_small")

# Function to extract text from file
def extract_text_from_file(filepath):
    if filepath.endswith(".txt"):
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif filepath.endswith(".pdf"):
        doc = fitz.open(filepath)
        return "\n".join([page.get_text() for page in doc])
    elif filepath.endswith(".pptx"):
        prs = Presentation(filepath)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    else:
        return ""

# Store raw documents in ChromaDB
raw_client = chromadb.PersistentClient(path=raw_doc_db_path)
raw_collection = raw_client.get_or_create_collection("raw_documents")

for filename in os.listdir(DOC_DIR):
    filepath = os.path.join(DOC_DIR, filename)
    if not os.path.isfile(filepath):
        continue

    text = extract_text_from_file(filepath)
    if not text.strip():
        print(f"Skipping empty or unsupported file: {filename}")
        continue

    try:
        raw_collection.add(documents=[text], ids=[filename])
        print(f"Stored raw document: {filename}")
    except Exception as e:
        print(f"Failed to store {filename}: {e}")

# Load stored raw documents
docs = raw_collection.get(include=["documents", "metadatas"])
texts = docs["documents"]
metadatas = docs["metadatas"]
ids = docs["ids"]

# Load embedding models
model_minilm = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
model_e5 = SentenceTransformer("intfloat/e5-small-v2")

# Helper to create & store embeddings
def store_embeddings(embedding_model, collection_name, client):
    collection = client.get_or_create_collection(name=collection_name, metadata={"hnsw:space": "cosine"})
    batch_size = 16
    existing_ids = set(collection.get()["ids"])  # get all existing IDs in one go

    added = 0
    for i in range(0, len(texts), batch_size):
        batch_texts = texts[i:i+batch_size]
        batch_ids = ids[i:i+batch_size]
        batch_metas = metadatas[i:i+batch_size]

        # Filter out existing IDs
        new_data = [
            (doc, emb_id, meta)
            for doc, emb_id, meta in zip(batch_texts, batch_ids, batch_metas)
            if emb_id not in existing_ids
        ]

        if not new_data:
            continue

        new_texts, new_ids, new_metas = zip(*new_data)
        new_texts = list(new_texts)
        new_ids = list(new_ids)
        new_metas = list(new_metas)

        embeddings = embedding_model.encode(new_texts, convert_to_numpy=True)
        embeddings = embeddings.tolist()  # Ensure list type for ChromaDB

        # Check lengths before adding
        if not (len(new_texts) == len(new_ids) == len(new_metas) == len(embeddings)):
            raise ValueError(f"Length mismatch before add: texts={len(new_texts)}, ids={len(new_ids)}, metas={len(new_metas)}, embeddings={len(embeddings)}")

        collection.add(documents=new_texts, embeddings=embeddings, metadatas=new_metas, ids=new_ids)
        added += len(new_ids)

    print(f"Saved {added} new embeddings to {collection_name}")


# Generate and store embeddings
minilm_client = chromadb.PersistentClient(path=embedding_db_minilm)
e5_client = chromadb.PersistentClient(path=embedding_db_e5)

store_embeddings(model_minilm, "embeddings_minilm", minilm_client)
store_embeddings(model_e5, "embeddings_e5_small", e5_client)
